# app.py
# Streamlit Healthcare Compliance Dashboard
# Author: Auto-generated for exam case study
#
# Features:
# - Weekly waiting time compliance trend
# - Department-wise compliance with traffic-light indicators
# - Doctor-level compliance + licensing risk detection
# - Export of tables & auto-generated PowerPoint summary
# - Flexible column mapping to fit *any* Excel schema
#
# How to run:
#   1) pip install -r requirements.txt
#   2) streamlit run app.py

import io
import os
import datetime as dt
from pathlib import Path

import numpy as np
import pandas as pd
import streamlit as st

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

APP_TITLE = "Healthcare Waiting-Time Compliance Dashboard"
DEFAULT_TARGET_MIN = 30     # default compliance target (minutes)
DEFAULT_ALERT_DAYS = 30     # license expiring soon threshold

st.set_page_config(page_title=APP_TITLE, page_icon="🩺", layout="wide")

# ---------- Helpers ----------
def load_data(file) -> pd.DataFrame:
    # Accept Excel or CSV
    if file is None:
        st.warning("Please upload your dataset (Excel/CSV) from the sidebar or keep the bundled example.")
        return pd.DataFrame()
    name = getattr(file, "name", "uploaded")
    try:
        if name.lower().endswith((".xlsx", ".xls")):
            df = pd.read_excel(file)
        elif name.lower().endswith(".csv"):
            df = pd.read_csv(file)
        else:
            # try Excel by default
            df = pd.read_excel(file)
    except Exception as e:
        st.error(f"Failed to read file: {e}")
        return pd.DataFrame()
    return df

def coerce_datetime(s):
    return pd.to_datetime(s, errors="coerce")

def coerce_numeric(s):
    return pd.to_numeric(s, errors="coerce")

def infer_week_start(d):
    # Monday as week start
    return (d - pd.to_timedelta((d.dt.dayofweek), unit="d")).dt.normalize()

def format_pct(x):
    if pd.isna(x):
        return ""
    return f"{x:.1f}%"

def traffic_light(p):
    # p in 0-100
    if pd.isna(p):
        return "⚪"
    return "🟢" if p >= dept_threshold else "🔴"

def figure_to_bytes(fig, dpi=150):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf

def make_powerpoint(summary, weekly_df, dept_df, doc_df, outfile):
    prs = Presentation()
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = APP_TITLE
    slide.placeholders[1].text = f"Generated on {dt.datetime.now():%Y-%m-%d %H:%M}\nDashboard Summary"
    
    # KPI Slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Key Findings"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for line in summary:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0

    # Charts Slide (Weekly + Department)
    chart_slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    chart_slide.shapes.title.text = "Weekly & Department Compliance"

    # Weekly line chart
    fig1, ax1 = plt.subplots(figsize=(6,3))
    if not weekly_df.empty:
        ax1.plot(weekly_df["week_start"], weekly_df["compliance_pct"])
        ax1.set_title("Weekly Waiting-Time Compliance (%)")
        ax1.set_xlabel("Week Start")
        ax1.set_ylabel("Compliance %")
        ax1.grid(True, alpha=0.3)
    img1 = figure_to_bytes(fig1)
    left = Inches(0.5); top = Inches(1.5); width = Inches(4.5)
    chart_slide.shapes.add_picture(img1, left, top, width=width)

    # Department bar chart
    fig2, ax2 = plt.subplots(figsize=(6,3))
    if not dept_df.empty:
        ax2.barh(dept_df["Department"], dept_df["Compliance %"])
        ax2.set_title("Department Compliance (%)")
        ax2.set_xlabel("Compliance %")
        ax2.set_ylabel("Department")
        ax2.grid(True, axis="x", alpha=0.3)
    img2 = figure_to_bytes(fig2)
    left = Inches(5.2); top = Inches(1.5); width = Inches(4.5)
    chart_slide.shapes.add_picture(img2, left, top, width=width)

    # Doctor risks slide (table-like text)
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Doctor Licensing Risks"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    if doc_df.empty:
        p = tf.add_paragraph(); p.text = "No licensing risks detected."; p.level = 0
    else:
        for _, r in doc_df.iterrows():
            p = tf.add_paragraph()
            p.text = f"{r['Doctor']} — License Expires: {r['License Expiry']} — Status: {r['Risk']}"
            p.level = 0

    prs.save(outfile)

# ---------- Sidebar: Data & Parameters ----------
st.sidebar.header("1) Data Source")
default_path = Path("data/Healthcare CaseStudy Data.xlsx")
use_bundled = st.sidebar.checkbox("Use bundled example", value=default_path.exists())
uploaded = None
if use_bundled and default_path.exists():
    uploaded = default_path.open("rb")
    uploaded.name = default_path.name
else:
    uploaded = st.sidebar.file_uploader("Upload Excel/CSV", type=["xlsx","xls","csv"])

raw = load_data(uploaded)
if raw.empty:
    st.stop()

st.sidebar.header("2) Map Columns")
cols = raw.columns.tolist()
# Common columns (choose the best match)
visit_date_col = st.sidebar.selectbox("Visit date column", cols, index=next((i for i,c in enumerate(cols) if 'date' in c.lower()), 0))
department_col = st.sidebar.selectbox("Department column", cols, index=next((i for i,c in enumerate(cols) if 'dept' in c.lower() or 'department' in c.lower()), 0))

# Waiting time: either a numeric column OR compute from two datetime columns
st.sidebar.subheader("Waiting Time")
waiting_mode = st.sidebar.radio("Compute waiting time from:", ["Numeric minutes column","Arrival & Seen times"])

if waiting_mode == "Numeric minutes column":
    waiting_minutes_col = st.sidebar.selectbox("Waiting minutes column", cols, index=next((i for i,c in enumerate(cols) if 'wait' in c.lower() and 'min' in c.lower()), 0))
    start_time_col = None; seen_time_col = None
else:
    start_time_col = st.sidebar.selectbox("Arrival/start time column", cols, index=next((i for i,c in enumerate(cols) if 'arrival' in c.lower() or 'check' in c.lower()), 0))
    seen_time_col = st.sidebar.selectbox("Seen/doctor start time column", cols, index=next((i for i,c in enumerate(cols) if 'seen' in c.lower() or 'start' in c.lower()), 0))
    waiting_minutes_col = None

doctor_col = st.sidebar.selectbox("Doctor column", cols, index=next((i for i,c in enumerate(cols) if 'doctor' in c.lower()), 0))

license_expiry_col = st.sidebar.selectbox("License expiry date column (optional)", ["<None>"] + cols, index=0)
license_expiry_col = None if license_expiry_col == "<None>" else license_expiry_col

st.sidebar.header("3) Parameters")
target_minutes = st.sidebar.number_input("Compliance target (max waiting minutes)", min_value=1, max_value=600, value=DEFAULT_TARGET_MIN, step=1)
dept_threshold = st.sidebar.slider("Department compliant if ≥ this %", min_value=0, max_value=100, value=90, step=1)
alert_days = st.sidebar.slider("License expiring soon (days)", min_value=1, max_value=365, value=DEFAULT_ALERT_DAYS, step=1)

# ---------- Prepare Data ----------
df = raw.copy()
# Coerce dates
df[visit_date_col] = coerce_datetime(df[visit_date_col])

if waiting_mode == "Numeric minutes column":
    df["waiting_minutes"] = coerce_numeric(df[waiting_minutes_col])
else:
    df[start_time_col] = coerce_datetime(df[start_time_col])
    df[seen_time_col] = coerce_datetime(df[seen_time_col])
    df["waiting_minutes"] = (df[seen_time_col] - df[start_time_col]).dt.total_seconds() / 60.0

df["is_compliant"] = df["waiting_minutes"] <= target_minutes

# Week start (Monday)
df = df.dropna(subset=[visit_date_col])
df["week_start"] = infer_week_start(df[visit_date_col])

# ---------- KPIs ----------
total_visits = len(df)
noncompliant_pct = (1.0 - df["is_compliant"].mean()) * 100 if total_visits > 0 else 0.0

colA, colB, colC, colD = st.columns(4)
colA.metric("Total Visits", f"{total_visits:,}")
colB.metric("% Noncompliant", f"{noncompliant_pct:.1f}%")
# Best/Worst departments by compliance
dept_stats = df.groupby(department_col)["is_compliant"].mean().mul(100).sort_values(ascending=False) if total_visits else pd.Series(dtype=float)
best_dept = dept_stats.index[0] if not dept_stats.empty else "-"
worst_dept = dept_stats.index[-1] if not dept_stats.empty else "-"
colC.metric("Best Dept", best_dept if isinstance(best_dept, str) else str(best_dept))
colD.metric("Worst Dept", worst_dept if isinstance(worst_dept, str) else str(worst_dept))

st.title(APP_TITLE)
st.caption("Upload any Excel/CSV with matching fields, map the columns in the sidebar, and the dashboard will adapt automatically.")

tab1, tab2, tab3, tab4, tab5 = st.tabs([
    "Weekly Trend", "Department Performance", "Doctors & Licensing", "Automation & Governance", "Export"
])

with tab1:
    st.subheader("Weekly Waiting-Time Compliance (%)")
    weekly = df.groupby("week_start")["is_compliant"].mean().mul(100).reset_index().rename(columns={"is_compliant":"compliance_pct"})
    st.line_chart(weekly.set_index("week_start"))
    st.dataframe(weekly)

with tab2:
    st.subheader("Department-wise Performance (Traffic Lights)")
    dept = df.groupby(department_col)["is_compliant"].mean().mul(100).reset_index()
    dept.columns = ["Department", "Compliance %"]
    dept["Indicator"] = dept["Compliance %"].apply(lambda p: "🟢" if p >= dept_threshold else "🔴")
    st.dataframe(dept.sort_values("Compliance %", ascending=False))
    st.bar_chart(dept.set_index("Department")["Compliance %"])

with tab3:
    st.subheader("Doctor-level Compliance & Licensing Issues")
    doctor = df.groupby(doctor_col)["is_compliant"].agg(["mean","count"]).reset_index()
    doctor["Compliance %"] = doctor["mean"]*100
    doctor.rename(columns={doctor_col:"Doctor", "count":"Visits"}, inplace=True)
    # Licensing
    today = pd.Timestamp.today().normalize()
    if license_expiry_col:
        df[license_expiry_col] = coerce_datetime(df[license_expiry_col])
        latest_expiry = df.groupby(doctor_col)[license_expiry_col].max().reset_index()
        latest_expiry.rename(columns={doctor_col:"Doctor", license_expiry_col:"License Expiry"}, inplace=True)
        doctor = doctor.merge(latest_expiry, on="Doctor", how="left")
        doctor["Days to Expiry"] = (doctor["License Expiry"] - today).dt.days
        def risk_label(row):
            d = row["Days to Expiry"]
            if pd.isna(d): return "Unknown"
            if d < 0: return "⛔ Expired"
            if d <= alert_days: return "⚠️ Expiring Soon"
            return "OK"
        doctor["Risk"] = doctor.apply(risk_label, axis=1)
    else:
        doctor["License Expiry"] = pd.NaT
        doctor["Days to Expiry"] = np.nan
        doctor["Risk"] = "Unknown"

    out_cols = ["Doctor", "Visits", "Compliance %", "License Expiry", "Days to Expiry", "Risk"]
    st.dataframe(doctor[out_cols].sort_values(["Risk","Compliance %"], ascending=[True, False]))

    risky_only = doctor[doctor["Risk"].isin(["⛔ Expired","⚠️ Expiring Soon"])][["Doctor","License Expiry","Days to Expiry","Risk"]]
    with st.expander("Show only at-risk doctors"):
        st.dataframe(risky_only)

with tab4:
    st.subheader("Automation & Governance – Practical Guide")
    st.markdown(f"""
**Goal:** refresh this dashboard daily and alert when performance drops.

**Option A – Python script + Task Scheduler/cron**
- Use the included `refresh_pipeline.py` to compute summarized CSVs.
- Schedule it daily:
  - **Windows Task Scheduler:** run `python refresh_pipeline.py` at 06:00.
  - **Linux/Mac cron:** `0 6 * * * /usr/bin/python3 /path/to/refresh_pipeline.py`.
- Point the Streamlit app to the refreshed CSVs or the source database.

**Option B – Direct DB connection**
- Replace the Excel reader with a database query (e.g., MSSQL, PostgreSQL).
- Run queries for yesterday/today and append to a central table.

**Alerts (Email/Teams/Slack)**
- After each refresh, if weekly compliance < **{dept_threshold}%** or % noncompliant > **{noncompliant_pct:.1f}%**, send an alert.
- Use SMTP or a webhook to notify Management & MoH auditors.

**Data governance**
- Define Owners: Data (Hospital IT), Process (Operations), Reporting (Compliance).
- Keep a data dictionary & change log.
- Access control with audit logs.
    """)

with tab5:
    st.subheader("Download Tables")
    # Weekly CSV
    weekly_csv = weekly.to_csv(index=False).encode("utf-8")
    st.download_button("Download Weekly Compliance (CSV)", weekly_csv, "weekly_compliance.csv", "text/csv")

    # Department CSV
    dept_csv = dept.to_csv(index=False).encode("utf-8")
    st.download_button("Download Department Performance (CSV)", dept_csv, "department_performance.csv", "text/csv")

    # Doctor CSV
    doctor_csv = doctor[out_cols].to_csv(index=False).encode("utf-8")
    st.download_button("Download Doctor Compliance & Licensing (CSV)", doctor_csv, "doctor_compliance_licensing.csv", "text/csv")

    st.divider()
    st.subheader("Generate PowerPoint (auto-summary)")
    if st.button("Create PPTX"):
        # Build summary bullets
        bullets = [
            f"Total visits analyzed: {total_visits:,}",
            f"Noncompliant visits: {noncompliant_pct:.1f}%",
        ]
        if not dept_stats.empty:
            bullets.append(f"Best department: {best_dept} ({dept_stats.iloc[0]:.1f}%)")
            bullets.append(f"Worst department: {worst_dept} ({dept_stats.iloc[-1]:.1f}%)")
        if license_expiry_col:
            exp_count = int((doctor['Risk'] == '⛔ Expired').sum())
            soon_count = int((doctor['Risk'] == '⚠️ Expiring Soon').sum())
            bullets.append(f"Licensing risks — Expired: {exp_count}, Expiring soon (≤{alert_days} days): {soon_count}")
        # Create PPT in-memory
        out_buf = io.BytesIO()
        make_powerpoint(bullets, weekly, dept, doctor[["Doctor","License Expiry","Risk"]], out_buf)
        st.download_button("Download Summary PPT", out_buf.getvalue(), file_name="Compliance_Summary.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

st.caption("Tip: Use the sidebar to fine-tune mappings/thresholds so this works with *any* hospital dataset.")
